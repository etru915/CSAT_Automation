from pptx import Presentation
import pandas as pd

prs = Presentation("C:\\PythonProjects\\ppt_word_checker\\inspection\\1 로켓직구 배송 심화.pptx")
file_name = "joy_CS"

text_runs = []
page_runs = []

page_count = 0
for slide in prs.slides:
    page_count += 1
    for shape in slide.shapes:
        # print(shape)
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text_runs.append(paragraph.text)
            page_runs.append("page " + str(page_count))
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        for n in shape.table.iter_cells():
            print(n.text)

        cell = shape.table.cell(0,0)
        # print(cell.text)
# print(text_runs)
# print(page_runs)

# df = pd.DataFrame()
# df["page_number"] = page_runs
# df["text_value"] = text_runs
# df["File_name"] = file_name
# df = df[["File_name","page_number","text_value"]]
#
# df.to_excel("C:\\CS Project\\Joy_CS.xlsx", encoding="EUC-KR")

# text_runs = []
# for slide in presentation.slides:
#     for shape in slide.shapes:
#         for run in shape.text_frame.paragraphs:
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     text_runs.append(run.text)
# #
# print(text_runs)